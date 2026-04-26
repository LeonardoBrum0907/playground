from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


TITLE_BLUE = RGBColor(17, 63, 115)
TEXT_DARK = RGBColor(31, 41, 55)
TEXT_MUTED = RGBColor(75, 85, 99)
BG_LIGHT = RGBColor(245, 247, 250)
CARD_BG = RGBColor(255, 255, 255)
BLUE_LIGHT = RGBColor(232, 242, 255)
GRAY_BORDER = RGBColor(210, 214, 220)
GREEN = RGBColor(46, 125, 50)
YELLOW = RGBColor(245, 158, 11)


def set_run_style(run, size, bold=False, color=TEXT_DARK):
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=14,
    bold=False,
    color=TEXT_DARK,
    align=PP_ALIGN.LEFT,
    fill=None,
    margin=0.08,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_style(run, size, bold=bold, color=color)
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.color.rgb = fill
    else:
        box.fill.background()
    return box


def add_bullet_textbox(slide, left, top, width, height, bullets, size=12):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.bullet = True
        p.space_after = Pt(8)
        for run in p.runs:
            set_run_style(run, size, color=TEXT_DARK)
    return box


def add_header(slide, title, subtitle):
    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.0)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = TITLE_BLUE
    banner.line.color.rgb = TITLE_BLUE
    add_textbox(
        slide,
        Inches(0.45),
        Inches(0.16),
        Inches(8.7),
        Inches(0.32),
        title,
        size=24,
        bold=True,
        color=RGBColor(255, 255, 255),
    )
    add_textbox(
        slide,
        Inches(0.45),
        Inches(0.52),
        Inches(8.7),
        Inches(0.20),
        subtitle,
        size=11,
        color=RGBColor(230, 238, 248),
    )
    add_textbox(
        slide,
        Inches(9.2),
        Inches(0.18),
        Inches(3.6),
        Inches(0.40),
        "Edite nome, gestor e periodo antes da apresentacao.",
        size=10,
        color=RGBColor(230, 238, 248),
        align=PP_ALIGN.RIGHT,
    )


def add_card(slide, left, top, width, height, title):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = GRAY_BORDER
    add_textbox(
        slide,
        left + Inches(0.16),
        top + Inches(0.10),
        width - Inches(0.32),
        Inches(0.24),
        title,
        size=15,
        bold=True,
        color=TITLE_BLUE,
    )
    return shape


def add_info_line(slide, text):
    add_textbox(
        slide,
        Inches(0.45),
        Inches(1.15),
        Inches(12.2),
        Inches(0.26),
        text,
        size=11,
        color=TEXT_MUTED,
    )


def add_action_row(slide, top, action, how, term):
    add_textbox(
        slide,
        Inches(0.70),
        top,
        Inches(2.50),
        Inches(0.52),
        action,
        size=11,
        bold=True,
        color=TEXT_DARK,
        fill=RGBColor(245, 248, 252),
        margin=0.05,
    )
    add_textbox(
        slide,
        Inches(3.30),
        top,
        Inches(6.55),
        Inches(0.52),
        how,
        size=10,
        color=TEXT_DARK,
        fill=RGBColor(252, 253, 255),
        margin=0.05,
    )
    add_textbox(
        slide,
        Inches(9.98),
        top,
        Inches(1.55),
        Inches(0.52),
        term,
        size=10,
        bold=True,
        color=TITLE_BLUE,
        align=PP_ALIGN.CENTER,
        fill=RGBColor(238, 245, 255),
        margin=0.05,
    )


def add_status_box(slide, left, top, width, title, items, color):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(1.72)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = GRAY_BORDER
    strip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + Inches(0.16), top + Inches(0.12), Inches(0.18), Inches(1.40)
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = color
    strip.line.color.rgb = color
    add_textbox(
        slide,
        left + Inches(0.45),
        top + Inches(0.12),
        width - Inches(0.60),
        Inches(0.24),
        title,
        size=13,
        bold=True,
        color=TITLE_BLUE,
    )
    add_bullet_textbox(
        slide,
        left + Inches(0.42),
        top + Inches(0.42),
        width - Inches(0.56),
        Inches(1.05),
        items,
        size=11,
    )


def build_presentation(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide1 = prs.slides.add_slide(blank)
    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = BG_LIGHT
    add_header(
        slide1,
        "PDI - Supervisor com Justificativas Excessivas",
        "Plano resumido com foco em escuta, objetividade e receptividade ao feedback.",
    )
    add_info_line(
        slide1,
        "Supervisor: [nome]   |   Gestor: [nome]   |   Periodo: [mes/ano a mes/ano]",
    )
    add_card(slide1, Inches(0.45), Inches(1.55), Inches(6.10), Inches(2.10), "Objetivo de desenvolvimento")
    add_textbox(
        slide1,
        Inches(0.65),
        Inches(2.00),
        Inches(5.70),
        Inches(1.35),
        "Desenvolver uma postura mais aberta ao feedback, reduzindo a necessidade de justificativas longas e repetitivas, com maior foco em escuta, objetividade, aprendizado e ajuste de conduta.",
        size=13,
        color=TEXT_DARK,
    )
    add_card(slide1, Inches(6.78), Inches(1.55), Inches(6.10), Inches(2.10), "Comportamento observado")
    add_textbox(
        slide1,
        Inches(6.98),
        Inches(2.00),
        Inches(5.70),
        Inches(1.35),
        "Em conversas de alinhamento, feedback ou correcao de rota, o supervisor tende a explicar excessivamente suas acoes, o que pode transmitir resistencia ao feedback e reduzir a objetividade da conversa.",
        size=13,
        color=TEXT_DARK,
    )
    add_card(slide1, Inches(0.45), Inches(3.90), Inches(12.43), Inches(2.55), "Resultado esperado")
    add_bullet_textbox(
        slide1,
        Inches(0.70),
        Inches(4.32),
        Inches(11.90),
        Inches(1.85),
        [
            "Ouvir o feedback ate o fim sem interromper ou entrar em defesa imediata.",
            "Explicar apenas o necessario, com objetividade e foco no contexto.",
            "Demonstrar abertura para ajustar a conduta sem insistir em justificativas repetitivas.",
            "Encerrar a conversa com um plano claro de acao e melhoria.",
        ],
        size=12,
    )

    slide2 = prs.slides.add_slide(blank)
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = BG_LIGHT
    add_header(
        slide2,
        "Acoes Rapidas do PDI",
        "Medidas simples para aplicar no dia a dia e acelerar a mudanca de comportamento.",
    )
    add_info_line(slide2, "Sugestao de acompanhamento: revisao semanal curta entre gestor e supervisor.")
    add_card(slide2, Inches(0.45), Inches(1.55), Inches(11.35), Inches(4.75), "Plano de acao resumido")
    add_textbox(slide2, Inches(0.70), Inches(1.98), Inches(2.40), Inches(0.22), "Acao", size=10, bold=True, color=TEXT_MUTED)
    add_textbox(slide2, Inches(3.30), Inches(1.98), Inches(6.10), Inches(0.22), "Como aplicar", size=10, bold=True, color=TEXT_MUTED)
    add_textbox(slide2, Inches(9.98), Inches(1.98), Inches(1.55), Inches(0.22), "Prazo", size=10, bold=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_action_row(
        slide2,
        Inches(2.28),
        "Pausa antes de responder",
        "Ao receber um feedback, ouvir o ponto completo e esperar alguns segundos antes de responder.",
        "Imediato",
    )
    add_action_row(
        slide2,
        Inches(2.92),
        "Resposta curta e objetiva",
        'Usar frases como "Entendi o ponto" ou "Vou ajustar essa conduta" antes de qualquer explicacao complementar.',
        "Imediato",
    )
    add_action_row(
        slide2,
        Inches(3.56),
        "Foco em melhoria",
        "Ao final de cada conversa, definir um ajuste pratico em vez de prolongar a justificativa.",
        "Semanal",
    )
    add_action_row(
        slide2,
        Inches(4.20),
        "Autoavaliacao rapida",
        "Registrar uma situacao por semana em que conseguiu ouvir melhor e responder com mais maturidade.",
        "Semanal",
    )
    add_card(slide2, Inches(11.98), Inches(1.55), Inches(0.90), Inches(4.75), "")
    add_textbox(
        slide2,
        Inches(12.02),
        Inches(2.00),
        Inches(0.82),
        Inches(3.90),
        "Dica:\nmenos defesa,\nmais escuta\ne acao.",
        size=12,
        bold=True,
        color=TITLE_BLUE,
        align=PP_ALIGN.CENTER,
        fill=BLUE_LIGHT,
        margin=0.08,
    )
    add_card(slide2, Inches(0.45), Inches(6.40), Inches(12.43), Inches(0.68), "Fala orientativa do gestor")
    add_textbox(
        slide2,
        Inches(0.62),
        Inches(6.67),
        Inches(12.05),
        Inches(0.28),
        '"Para evoluir, preciso que voce reduza a necessidade de se justificar em excesso, ouca o feedback com mais abertura e transforme a conversa em ajuste pratico de conduta."',
        size=11,
        color=TEXT_DARK,
    )

    slide3 = prs.slides.add_slide(blank)
    slide3.background.fill.solid()
    slide3.background.fill.fore_color.rgb = BG_LIGHT
    add_header(
        slide3,
        "Acompanhamento e Evidencias",
        "Indicadores simples para mostrar evolucao de forma objetiva.",
    )
    add_info_line(slide3, "Preencha com exemplos reais observados durante o periodo do PDI.")
    add_status_box(
        slide3,
        Inches(0.45),
        Inches(1.55),
        Inches(4.00),
        "Evidencias esperadas",
        [
            "Recebe feedback sem interromper.",
            "Reduz justificativas longas.",
            "Aceita redirecionamentos com menor defensividade.",
            "Foca em ajuste e nao em convencimento.",
        ],
        GREEN,
    )
    add_status_box(
        slide3,
        Inches(4.67),
        Inches(1.55),
        Inches(4.00),
        "Como registrar a evolucao",
        [
            "Anotar situacao, reacao e resultado.",
            "Comparar antes e depois do acompanhamento.",
            "Coletar percepcao do gestor e da equipe.",
            "Usar revisoes semanais curtas.",
        ],
        YELLOW,
    )
    add_card(slide3, Inches(8.90), Inches(1.55), Inches(3.98), Inches(2.65), "Quadro de acompanhamento")
    add_textbox(
        slide3,
        Inches(9.12),
        Inches(2.00),
        Inches(3.52),
        Inches(1.85),
        "Semana 1: ____________________\n\nSemana 2: ____________________\n\nSemana 3: ____________________\n\nSemana 4: ____________________",
        size=12,
        color=TEXT_DARK,
        fill=RGBColor(252, 253, 255),
        margin=0.08,
    )
    add_card(slide3, Inches(0.45), Inches(4.50), Inches(12.43), Inches(2.02), "Leitura final para apresentacao")
    add_textbox(
        slide3,
        Inches(0.68),
        Inches(4.95),
        Inches(11.95),
        Inches(1.20),
        "A evolucao sera percebida quando o supervisor passar a ouvir o feedback com mais abertura, explicar apenas o necessario e demonstrar mudanca por meio de acoes praticas. O foco da avaliacao deve estar no comportamento observavel e nao apenas no discurso.",
        size=13,
        color=TEXT_DARK,
    )

    prs.save(output_path)


if __name__ == "__main__":
    build_presentation("/workspace/PDI_supervisor_justificativas.pptx")
