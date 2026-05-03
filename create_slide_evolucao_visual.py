from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


IMG_PATH = "/workspace/imagem_evolucao_pdi_visual.png"
PPT_PATH = "/workspace/slide_evolucao_evidencias.pptx"

TITLE_BLUE = RGBColor(17, 63, 115)
TEXT_DARK = RGBColor(31, 41, 55)
TEXT_MUTED = RGBColor(75, 85, 99)
BG_LIGHT = RGBColor(245, 247, 250)
CARD_BG = RGBColor(255, 255, 255)
GRAY_BORDER = RGBColor(210, 214, 220)
GREEN = RGBColor(52, 168, 83)
YELLOW = RGBColor(245, 158, 11)
RED = RGBColor(220, 38, 38)
BLUE_LIGHT = RGBColor(232, 242, 255)


def set_run_style(run, size, bold=False, color=TEXT_DARK):
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=14,
    bold=False,
    color=TEXT_DARK,
    align=PP_ALIGN.LEFT,
    fill=None,
    margin=0.08,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_style(run, size, bold=bold, color=color)
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.color.rgb = fill
    else:
        box.fill.background()
    return box


def add_bullets(slide, left, top, width, height, title, bullets):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = GRAY_BORDER

    add_textbox(
        slide,
        left + Inches(0.16),
        top + Inches(0.10),
        width - Inches(0.32),
        Inches(0.25),
        title,
        size=14,
        bold=True,
        color=TITLE_BLUE,
    )

    box = slide.shapes.add_textbox(
        left + Inches(0.16), top + Inches(0.42), width - Inches(0.32), height - Inches(0.52)
    )
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.bullet = True
        p.space_after = Pt(8)
        for run in p.runs:
            set_run_style(run, 11, color=TEXT_DARK)


def draw_illustration(path):
    image = Image.new("RGB", (1600, 900), (245, 247, 250))
    draw = ImageDraw.Draw(image)

    # Background cards
    draw.rounded_rectangle((55, 70, 1545, 830), radius=36, fill=(255, 255, 255), outline=(220, 226, 232), width=3)
    draw.rounded_rectangle((90, 110, 760, 785), radius=28, fill=(236, 244, 255), outline=(201, 220, 245), width=2)
    draw.rounded_rectangle((820, 110, 1510, 785), radius=28, fill=(248, 250, 252), outline=(225, 230, 236), width=2)

    # Left panel: person icon
    draw.ellipse((230, 220, 380, 370), fill=(17, 63, 115))
    draw.rounded_rectangle((180, 370, 430, 670), radius=36, fill=(52, 97, 151))
    draw.rounded_rectangle((208, 450, 402, 580), radius=24, fill=(232, 242, 255))

    # Speech bubbles representing communication
    draw.rounded_rectangle((450, 210, 650, 300), radius=22, fill=(255, 255, 255), outline=(201, 220, 245), width=3)
    draw.polygon([(470, 300), (495, 300), (478, 325)], fill=(255, 255, 255), outline=(201, 220, 245))
    draw.rounded_rectangle((430, 360, 690, 455), radius=22, fill=(255, 255, 255), outline=(201, 220, 245), width=3)
    draw.polygon([(450, 455), (475, 455), (458, 482)], fill=(255, 255, 255), outline=(201, 220, 245))

    # Simple icons / bars
    draw.rectangle((475, 238, 505, 275), fill=(52, 168, 83))
    draw.rectangle((518, 225, 548, 275), fill=(245, 158, 11))
    draw.rectangle((561, 205, 591, 275), fill=(17, 63, 115))
    draw.line((470, 408, 510, 395, 550, 410, 600, 375, 650, 340), fill=(52, 168, 83), width=8)

    # Right panel chart illustration
    chart_x0, chart_y0 = 900, 540
    chart_w, chart_h = 500, 170
    draw.line((chart_x0, chart_y0 - chart_h, chart_x0, chart_y0), fill=(170, 180, 190), width=4)
    draw.line((chart_x0, chart_y0, chart_x0 + chart_w, chart_y0), fill=(170, 180, 190), width=4)

    bar_width = 80
    gaps = 40
    heights = [65, 105, 130, 155]
    colors = [(220, 38, 38), (245, 158, 11), (82, 130, 200), (52, 168, 83)]
    start_x = chart_x0 + 40
    for idx, (height, color) in enumerate(zip(heights, colors)):
        x0 = start_x + idx * (bar_width + gaps)
        y0 = chart_y0 - height
        x1 = x0 + bar_width
        draw.rounded_rectangle((x0, y0, x1, chart_y0), radius=14, fill=color)

    draw.line(
        (
            start_x + 40,
            chart_y0 - heights[0] - 20,
            start_x + 160,
            chart_y0 - heights[1] - 30,
            start_x + 280,
            chart_y0 - heights[2] - 28,
            start_x + 400,
            chart_y0 - heights[3] - 42,
        ),
        fill=(17, 63, 115),
        width=10,
    )
    draw.polygon(
        [
            (start_x + 400, chart_y0 - heights[3] - 42),
            (start_x + 430, chart_y0 - heights[3] - 52),
            (start_x + 410, chart_y0 - heights[3] - 15),
        ],
        fill=(17, 63, 115),
    )

    # Small signal cards
    for y, label_color in [(170, (52, 168, 83)), (250, (245, 158, 11)), (330, (17, 63, 115))]:
        draw.rounded_rectangle((910, y, 1460, y + 56), radius=18, fill=(255, 255, 255), outline=(225, 230, 236), width=2)
        draw.ellipse((935, y + 15, 965, y + 45), fill=label_color)
        draw.rounded_rectangle((985, y + 18, 1390, y + 38), radius=10, fill=(224, 231, 239))

    image.save(path)


def build_presentation(image_path, output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_LIGHT

    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.0)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = TITLE_BLUE
    banner.line.color.rgb = TITLE_BLUE

    add_textbox(
        slide,
        Inches(0.45),
        Inches(0.16),
        Inches(8.4),
        Inches(0.32),
        "Evolução e Evidências do PDI",
        size=24,
        bold=True,
        color=RGBColor(255, 255, 255),
    )
    add_textbox(
        slide,
        Inches(0.45),
        Inches(0.52),
        Inches(9.6),
        Inches(0.22),
        "Slide executivo com imagem ilustrativa, gráfico de evolução e evidências comportamentais.",
        size=11,
        color=RGBColor(230, 238, 248),
    )
    add_textbox(
        slide,
        Inches(9.25),
        Inches(0.18),
        Inches(3.6),
        Inches(0.22),
        "Edite os textos conforme o caso real.",
        size=10,
        color=RGBColor(230, 238, 248),
        align=PP_ALIGN.RIGHT,
    )

    add_textbox(
        slide,
        Inches(0.45),
        Inches(1.15),
        Inches(7.4),
        Inches(0.20),
        "Supervisor: [nome]   |   Período: [mês/ano a mês/ano]",
        size=11,
        color=TEXT_MUTED,
    )

    # Left image panel
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(1.5), Inches(6.0), Inches(5.35)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = CARD_BG
    panel.line.color.rgb = GRAY_BORDER
    add_textbox(slide, Inches(0.62), Inches(1.63), Inches(4.2), Inches(0.24), "Imagem da evolução", size=14, bold=True, color=TITLE_BLUE)
    slide.shapes.add_picture(image_path, Inches(0.70), Inches(2.00), width=Inches(5.5), height=Inches(3.15))
    add_textbox(
        slide,
        Inches(0.70),
        Inches(5.30),
        Inches(5.5),
        Inches(1.25),
        "Representação visual da mudança de postura: mais escuta, mais equilíbrio e maior foco em solução.",
        size=12,
        color=TEXT_DARK,
        fill=BLUE_LIGHT,
    )

    # Right chart panel
    chart_card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.75), Inches(1.5), Inches(6.13), Inches(2.6)
    )
    chart_card.fill.solid()
    chart_card.fill.fore_color.rgb = CARD_BG
    chart_card.line.color.rgb = GRAY_BORDER
    add_textbox(slide, Inches(6.95), Inches(1.63), Inches(4.5), Inches(0.24), "Gráfico de evolução", size=14, bold=True, color=TITLE_BLUE)

    chart_data = CategoryChartData()
    chart_data.categories = ["Início", "Meio", "Atual"]
    chart_data.add_series("Controle emocional", (2, 3, 4))
    chart_data.add_series("Menos justificativa", (2, 3, 4))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(7.00),
        Inches(1.98),
        Inches(5.60),
        Inches(1.80),
        chart_data,
    ).chart
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 5
    chart.value_axis.major_unit = 1
    chart.has_legend = True
    chart.legend.position = 2
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = TITLE_BLUE
    chart.series[1].format.fill.solid()
    chart.series[1].format.fill.fore_color.rgb = GREEN

    add_bullets(
        slide,
        Inches(6.75),
        Inches(4.35),
        Inches(6.13),
        Inches(2.50),
        "Evidências para apresentar",
        [
            "Controle emocional: passou a reagir com mais equilíbrio diante de negativas e feedbacks.",
            "Controle emocional: demonstra mais escuta antes da resposta e menor impulsividade.",
            "Menos justificativa: reduziu explicações excessivas e tornou as respostas mais objetivas.",
            "Menos justificativa: mostra mais abertura para ajuste de conduta e menos postura defensiva.",
        ],
    )

    footer = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(6.95), Inches(12.43), Inches(0.35)
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = BLUE_LIGHT
    footer.line.color.rgb = RGBColor(198, 214, 241)
    add_textbox(
        slide,
        Inches(0.55),
        Inches(7.00),
        Inches(12.1),
        Inches(0.16),
        "Mensagem final: a evolução é percebida pela combinação entre postura mais equilibrada, maior objetividade e foco em solução.",
        size=10,
        color=TITLE_BLUE,
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0.03,
    )

    prs.save(output_path)


if __name__ == "__main__":
    draw_illustration(IMG_PATH)
    build_presentation(IMG_PATH, PPT_PATH)
