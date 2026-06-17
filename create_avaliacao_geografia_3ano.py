#!/usr/bin/env python3
"""Gera avaliação de Geografia para o 3º ano do Ensino Fundamental."""

import os
import time
import urllib.request
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT_PATH = "/workspace/Avaliacao_Geografia_3ano.pptx"
ASSETS_DIR = Path("/workspace/avaliacao_geografia_assets")

TITLE_GREEN = RGBColor(27, 94, 32)
TITLE_BLUE = RGBColor(17, 62, 115)
TEXT_DARK = RGBColor(31, 41, 55)
TEXT_MUTED = RGBColor(75, 85, 99)
BG_LIGHT = RGBColor(245, 247, 250)
CARD_BG = RGBColor(255, 255, 255)
GRAY_BORDER = RGBColor(210, 214, 220)
ACCENT = RGBColor(46, 125, 50)

IMAGE_SOURCES = {
    "bairro": (
        "https://upload.wikimedia.org/wikipedia/commons/thumb/a/ad/"
        "Andrel%C3%A2ndia_MG_Brasil_-_Bairro_dos_Maraj%C3%A1s.JPG/"
        "960px-Andrel%C3%A2ndia_MG_Brasil_-_Bairro_dos_Maraj%C3%A1s.JPG"
    ),
    "escola": (
        "https://upload.wikimedia.org/wikipedia/commons/thumb/5/54/"
        "Alunos_Escola_Municipal_Gyn_Brasil.jpg/"
        "960px-Alunos_Escola_Municipal_Gyn_Brasil.jpg"
    ),
    "cultura_festa": (
        "https://upload.wikimedia.org/wikipedia/commons/thumb/a/a3/"
        "Dan%C3%A7a_de_crian%C3%A7as_do_ensino_fundamental_1_na_festa_junina_do_"
        "Col%C3%A9gio_Padre_de_Man%2C_Coronel_Fabriciano_MG.JPG/"
        "960px-Dan%C3%A7a_de_crian%C3%A7as_do_ensino_fundamental_1_na_festa_junina_"
        "do_Col%C3%A9gio_Padre_de_Man%2C_Coronel_Fabriciano_MG.JPG"
    ),
    "geracoes": (
        "https://upload.wikimedia.org/wikipedia/commons/thumb/9/92/"
        "Three_generations_of_a_family_%281%29.jpg/"
        "960px-Three_generations_of_a_family_%281%29.jpg"
    ),
    "cultura_urbana": (
        "https://upload.wikimedia.org/wikipedia/commons/thumb/3/3b/"
        "Grafite_no_BankBoston_da_Avenida_Paulista%2C_S%C3%A3o_Paulo.jpg/"
        "960px-Grafite_no_BankBoston_da_Avenida_Paulista%2C_S%C3%A3o_Paulo.jpg"
    ),
    "samba_urbano": (
        "https://upload.wikimedia.org/wikipedia/commons/thumb/1/18/"
        "Blocos_de_rua_abrem_o_carnaval_de_S%C3%A3o_Paulo_%2832821515942%29.jpg/"
        "960px-Blocos_de_rua_abrem_o_carnaval_de_S%C3%A3o_Paulo_%2832821515942%29.jpg"
    ),
}

FALLBACK_LABELS = {
    "bairro": "Bairro residencial",
    "escola": "Escola municipal",
    "cultura_festa": "Festa junina",
    "geracoes": "Três gerações",
    "cultura_urbana": "Arte urbana",
    "samba_urbano": "Samba urbano",
}


def set_run_style(run, size, bold=False, color=TEXT_DARK):
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=14,
    bold=False,
    color=TEXT_DARK,
    align=PP_ALIGN.LEFT,
    fill=None,
    margin=0.08,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_style(run, size, bold=bold, color=color)
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.color.rgb = fill
    else:
        box.fill.background()
    return box


def add_header(slide, title, subtitle=""):
    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.95)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = TITLE_GREEN
    banner.line.color.rgb = TITLE_GREEN
    add_textbox(
        slide,
        Inches(0.45),
        Inches(0.18),
        Inches(12.4),
        Inches(0.35),
        title,
        size=22,
        bold=True,
        color=RGBColor(255, 255, 255),
    )
    if subtitle:
        add_textbox(
            slide,
            Inches(0.45),
            Inches(0.55),
            Inches(12.4),
            Inches(0.25),
            subtitle,
            size=11,
            color=RGBColor(230, 245, 232),
        )


def add_question_block(slide, number, theme, question, options=None, lines=3):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.45),
        Inches(1.05),
        Inches(7.55),
        Inches(6.15),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = GRAY_BORDER

    add_textbox(
        slide,
        Inches(0.65),
        Inches(1.18),
        Inches(7.15),
        Inches(0.22),
        f"Questão {number}  •  {theme}",
        size=10,
        bold=True,
        color=ACCENT,
    )
    add_textbox(
        slide,
        Inches(0.65),
        Inches(1.48),
        Inches(7.15),
        Inches(1.35),
        question,
        size=13,
        bold=True,
        color=TEXT_DARK,
    )

    y = 2.95
    if options:
        for opt in options:
            add_textbox(
                slide,
                Inches(0.75),
                Inches(y),
                Inches(7.0),
                Inches(0.42),
                opt,
                size=12,
                color=TEXT_DARK,
            )
            y += 0.48
    else:
        add_textbox(
            slide,
            Inches(0.75),
            Inches(y),
            Inches(7.0),
            Inches(0.25),
            "Resposta:",
            size=11,
            bold=True,
            color=TEXT_MUTED,
        )
        for i in range(lines):
            line_y = y + 0.35 + (i * 0.55)
            line = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(0.75),
                Inches(line_y),
                Inches(6.95),
                Inches(0.02),
            )
            line.fill.solid()
            line.fill.fore_color.rgb = GRAY_BORDER
            line.line.color.rgb = GRAY_BORDER


def add_image_panel(slide, image_path, caption):
    frame = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(8.25),
        Inches(1.05),
        Inches(4.65),
        Inches(6.15),
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = CARD_BG
    frame.line.color.rgb = GRAY_BORDER

    slide.shapes.add_picture(
        str(image_path),
        Inches(8.45),
        Inches(1.25),
        width=Inches(4.25),
        height=Inches(4.85),
    )
    add_textbox(
        slide,
        Inches(8.45),
        Inches(6.20),
        Inches(4.25),
        Inches(0.75),
        caption,
        size=10,
        color=TEXT_MUTED,
        align=PP_ALIGN.CENTER,
    )


def create_fallback_image(name, label):
    img = Image.new("RGB", (640, 480), (232, 245, 233))
    draw = ImageDraw.Draw(img)
    draw.rectangle([20, 20, 620, 460], outline=(46, 125, 50), width=4)
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 28)
    except OSError:
        font = ImageFont.load_default()
    draw.multiline_text(
        (60, 200),
        f"{label}\n(imagem ilustrativa)",
        fill=(27, 94, 32),
        font=font,
        align="center",
    )
    path = ASSETS_DIR / f"{name}.jpg"
    img.save(path, quality=90)
    return path


def download_images():
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)
    paths = {}
    headers = {"User-Agent": "GeoAssessmentBot/1.0 (educational; contact: classroom use)"}

    for name, url in IMAGE_SOURCES.items():
        path = ASSETS_DIR / f"{name}.jpg"
        downloaded = False
        for attempt in range(3):
            try:
                req = urllib.request.Request(url, headers=headers)
                with urllib.request.urlopen(req, timeout=30) as response:
                    data = response.read()
                if len(data) < 30000:
                    raise ValueError("Arquivo muito pequeno")
                path.write_bytes(data)
                paths[name] = path
                downloaded = True
                break
            except Exception:
                time.sleep(1.5 * (attempt + 1))
        if not downloaded:
            if path.exists() and path.stat().st_size >= 30000:
                paths[name] = path
            else:
                paths[name] = create_fallback_image(name, FALLBACK_LABELS[name])
        time.sleep(0.4)
    return paths


QUESTIONS = [
    {
        "number": 1,
        "theme": "Lugares de vivência",
        "question": (
            "Observe a imagem ao lado. Qual alternativa descreve melhor um "
            "lugar de vivência?"
        ),
        "options": [
            "a) Um local onde vivemos, convivemos e realizamos atividades do dia a dia.",
            "b) Apenas lugares turísticos visitados nas férias.",
            "c) Somente a escola, pois é onde aprendemos.",
            "d) Lugares que existem apenas em mapas e atlas.",
        ],
        "image": "bairro",
        "caption": "Bairro residencial — um dos lugares onde muitas pessoas vivem e convivem.",
        "answer": "a",
    },
    {
        "number": 2,
        "theme": "Lugares de vivência",
        "question": (
            "Mariana mora com a família, estuda de manhã, brinca na praça do "
            "bairro e vai ao mercado com a mãe. Quantos lugares de vivência "
            "diferentes aparecem nessa descrição?"
        ),
        "options": [
            "a) Apenas um, porque tudo acontece no mesmo bairro.",
            "b) Dois: casa e escola.",
            "c) Quatro ou mais: casa, escola, praça e mercado.",
            "d) Nenhum, pois lugares de vivência são só cidades grandes.",
        ],
        "image": "escola",
        "caption": "Escola — lugar de vivência onde aprendemos e convivemos com colegas.",
        "answer": "c",
    },
    {
        "number": 3,
        "theme": "Lugares de vivência",
        "question": (
            "Leia as afirmativas e marque V (verdadeiro) ou F (falso):\n\n"
            "(  ) Lugares de vivência incluem casa, escola e bairro.\n"
            "(  ) Só existem lugares de vivência nas capitais.\n"
            "(  ) Cada pessoa pode ter lugares de vivência diferentes."
        ),
        "options": None,
        "lines": 1,
        "image": "bairro",
        "caption": "Nosso bairro faz parte dos lugares onde construímos nossa rotina.",
        "answer": "V, F, V",
    },
    {
        "number": 4,
        "theme": "O que é cultura",
        "question": (
            "Cultura é o conjunto de costumes, crenças, festas, comidas, "
            "músicas e formas de viver de um grupo de pessoas. Com base nisso, "
            "qual item NÃO representa um elemento cultural?"
        ),
        "options": [
            "a) Preparar comidas típicas da região em datas especiais.",
            "b) Celebrar festas tradicionais com danças e músicas.",
            "c) A altura de uma montanha medida em metros.",
            "d) Contar histórias passadas de pais para filhos.",
        ],
        "image": "cultura_festa",
        "caption": "Festas juninas reúnem comidas, danças e tradições culturais.",
        "answer": "c",
    },
    {
        "number": 5,
        "theme": "O que é cultura",
        "question": (
            "Observe a imagem. A festa representada é um exemplo de manifestação "
            "cultural porque:"
        ),
        "options": [
            "a) Acontece somente em um único bairro e não pode ser compartilhada.",
            "b) Reúne tradições, danças, comidas e formas de celebrar de um povo.",
            "c) Não muda com o passar do tempo.",
            "d) Existe apenas para entreter turistas estrangeiros.",
        ],
        "image": "cultura_festa",
        "caption": "Manifestações culturais expressam identidade e pertencimento.",
        "answer": "b",
    },
    {
        "number": 6,
        "theme": "O que é cultura",
        "question": (
            "Complete a frase: A cultura de um grupo é formada por tudo aquilo "
            "que as pessoas __________ ao longo do tempo e transmitem umas "
            "às outras."
        ),
        "options": None,
        "lines": 2,
        "image": "cultura_festa",
        "caption": "Danças e trajes típicos fazem parte da cultura compartilhada.",
        "answer": "criam, praticam e transmitem (ou: aprendem e compartilham)",
    },
    {
        "number": 7,
        "theme": "Cultura entre gerações",
        "question": (
            "Observe a imagem com pessoas de idades diferentes. Por que a "
            "cultura pode ser diferente entre avós, pais e crianças?"
        ),
        "options": [
            "a) Porque cada geração vive em épocas e contextos distintos.",
            "b) Porque apenas os avós têm cultura de verdade.",
            "c) Porque crianças não participam de nenhuma tradição.",
            "d) Porque a cultura nunca muda de uma geração para outra.",
        ],
        "image": "geracoes",
        "caption": "Diferentes gerações convivem e compartilham tradições.",
        "answer": "a",
    },
    {
        "number": 8,
        "theme": "Cultura entre gerações",
        "question": (
            "Relacione as colunas:\n\n"
            "1. Avós          A. Brincadeiras de rua e piões\n"
            "2. Pais          B. Jogos digitais e vídeos na internet\n"
            "3. Crianças      C. Festas de família e receitas antigas\n\n"
            "Escreva a sequência correta (ex.: 1-C, 2-A, 3-B):"
        ),
        "options": None,
        "lines": 1,
        "image": "geracoes",
        "caption": "Cada geração tem manifestações culturais próprias e também compartilhadas.",
        "answer": "1-C, 2-A, 3-B",
    },
    {
        "number": 9,
        "theme": "Cultura entre gerações",
        "question": (
            "Qual alternativa apresenta duas manifestações culturais que podem "
            "ser compartilhadas entre gerações?"
        ),
        "options": [
            "a) Samba de roda e contação de histórias em família.",
            "b) Tipo de celular e marca do videogame.",
            "c) Cor do uniforme escolar e número da sala.",
            "d) Altura de um prédio e largura de uma rua.",
        ],
        "image": "samba_urbano",
        "caption": "Música e dança podem unir pessoas de diferentes idades.",
        "answer": "a",
    },
    {
        "number": 10,
        "theme": "Cultura brasileira urbana",
        "question": (
            "Observe a imagem de arte urbana. Esse tipo de manifestação "
            "cultural é comum em cidades brasileiras porque:"
        ),
        "options": [
            "a) Só existe em países frios, longe do Brasil.",
            "b) Expressa identidade, criatividade e vida nas cidades.",
            "c) Não tem relação com a cultura local.",
            "d) Substitui completamente todas as outras tradições.",
        ],
        "image": "cultura_urbana",
        "caption": "Grafites e murais fazem parte das culturas urbanas brasileiras.",
        "answer": "b",
    },
    {
        "number": 11,
        "theme": "Cultura brasileira urbana",
        "question": (
            "São exemplos de culturas urbanas no Brasil, EXCETO:"
        ),
        "options": [
            "a) Funk carioca e batidas das comunidades.",
            "b) Grafite e arte de rua em grandes cidades.",
            "c) Samba de roda e rodas de capoeira em praças.",
            "d) Formação de desertos de areia no interior do país.",
        ],
        "image": "cultura_urbana",
        "caption": "Nas cidades, a cultura se manifesta nas ruas, praças e bairros.",
        "answer": "d",
    },
    {
        "number": 12,
        "theme": "Cultura brasileira urbana",
        "question": (
            "Observe a imagem de um bloco de carnaval nas ruas de São Paulo. "
            "Essa manifestação mostra que a cultura brasileira urbana:"
        ),
        "options": [
            "a) É igual em todas as regiões, sem diferenças.",
            "b) Une música, dança, história e convivência no espaço da cidade.",
            "c) Só acontece em áreas rurais distantes.",
            "d) Não tem importância para a identidade do povo brasileiro.",
        ],
        "image": "samba_urbano",
        "caption": "Blocos de carnaval nas ruas são manifestações da cultura urbana brasileira.",
        "answer": "b",
    },
]

GABARITO = [
    ("1", "a", "Lugares de vivência"),
    ("2", "c", "Lugares de vivência"),
    ("3", "V, F, V", "Lugares de vivência"),
    ("4", "c", "O que é cultura"),
    ("5", "b", "O que é cultura"),
    ("6", "criam, praticam e transmitem", "O que é cultura"),
    ("7", "a", "Cultura entre gerações"),
    ("8", "1-C, 2-A, 3-B", "Cultura entre gerações"),
    ("9", "a", "Cultura entre gerações"),
    ("10", "b", "Cultura brasileira urbana"),
    ("11", "d", "Cultura brasileira urbana"),
    ("12", "b", "Cultura brasileira urbana"),
]


def build_cover(prs, blank):
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_LIGHT

    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = TITLE_GREEN
    banner.line.color.rgb = TITLE_GREEN

    add_textbox(
        slide,
        Inches(0.8),
        Inches(1.8),
        Inches(11.7),
        Inches(0.8),
        "Avaliação de Geografia",
        size=36,
        bold=True,
        color=RGBColor(255, 255, 255),
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(2.75),
        Inches(11.7),
        Inches(0.5),
        "3º Ano do Ensino Fundamental",
        size=22,
        color=RGBColor(230, 245, 232),
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(2.0),
        Inches(3.8),
        Inches(9.3),
        Inches(1.2),
        "Conteúdos: lugares de vivência • o que é cultura • cultura entre gerações • cultura brasileira urbana",
        size=14,
        color=RGBColor(255, 255, 255),
        align=PP_ALIGN.CENTER,
    )

    info = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(3.2),
        Inches(5.2),
        Inches(6.9),
        Inches(1.5),
    )
    info.fill.solid()
    info.fill.fore_color.rgb = CARD_BG
    info.line.color.rgb = GRAY_BORDER
    add_textbox(
        slide,
        Inches(3.45),
        Inches(5.45),
        Inches(6.4),
        Inches(1.0),
        "Aluno(a): ________________________________\n\nData: ____/____/________    Turma: __________",
        size=14,
        color=TEXT_DARK,
        align=PP_ALIGN.CENTER,
    )


def build_instructions(prs, blank):
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_LIGHT
    add_header(slide, "Instruções", "Leia com atenção antes de responder.")

    add_textbox(
        slide,
        Inches(0.65),
        Inches(1.25),
        Inches(12.0),
        Inches(5.5),
        (
            "• Esta avaliação contém 12 questões de nível médio sobre Geografia.\n\n"
            "• Observe atentamente as imagens ao lado de cada questão — elas ajudam na resposta.\n\n"
            "• Questões objetivas: marque apenas uma alternativa com X.\n\n"
            "• Questões discursivas: escreva com letra legível nas linhas indicadas.\n\n"
            "• Valor total: 12 pontos (1 ponto por questão).\n\n"
            "• Tempo sugerido: 50 minutos.\n\n"
            "• Boa prova!"
        ),
        size=15,
        color=TEXT_DARK,
    )


def build_question_slides(prs, blank, images):
    for q in QUESTIONS:
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG_LIGHT
        add_header(slide, "Avaliação de Geografia — 3º Ano", f"Questão {q['number']} de 12")
        add_question_block(
            slide,
            q["number"],
            q["theme"],
            q["question"],
            options=q.get("options"),
            lines=q.get("lines", 3),
        )
        add_image_panel(slide, images[q["image"]], q["caption"])


def build_answer_key(prs, blank):
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_LIGHT
    add_header(slide, "Gabarito — Uso do professor(a)", "Respostas esperadas")

    y = 1.25
    for num, answer, theme in GABARITO:
        add_textbox(
            slide,
            Inches(0.65),
            Inches(y),
            Inches(12.0),
            Inches(0.32),
            f"Questão {num} ({theme}): {answer}",
            size=12,
            color=TEXT_DARK,
        )
        y += 0.42

    add_textbox(
        slide,
        Inches(0.65),
        Inches(6.55),
        Inches(12.0),
        Inches(0.5),
        "Observação: na questão 6, aceitar respostas equivalentes que expressem criação, prática e transmissão cultural.",
        size=10,
        color=TEXT_MUTED,
    )


def build_presentation(output_path):
    images = download_images()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    build_cover(prs, blank)
    build_instructions(prs, blank)
    build_question_slides(prs, blank, images)
    build_answer_key(prs, blank)

    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    path = build_presentation(OUTPUT_PATH)
    print(f"Avaliação gerada: {path}")
