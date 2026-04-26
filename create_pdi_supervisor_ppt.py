from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


TITLE_BLUE = RGBColor(17, 62, 115)
TEXT_DARK = RGBColor(31, 41, 55)
TEXT_MUTED = RGBColor(75, 85, 99)
BG_LIGHT = RGBColor(245, 247, 250)
CARD_BG = RGBColor(255, 255, 255)
GREEN = RGBColor(46, 125, 50)
YELLOW = RGBColor(245, 158, 11)
RED = RGBColor(198, 40, 40)
BLUE_LIGHT = RGBColor(224, 239, 255)
GRAY_BORDER = RGBColor(210, 214, 220)


def set_run_style(run, size, bold=False, color=TEXT_DARK):
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=16, bold=False,
                color=TEXT_DARK, align=PP_ALIGN.LEFT, fill=None, margin=0.10):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_style(run, size, bold=bold, color=color)
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.color.rgb = fill
    else:
        box.fill.background()
    return box


def add_bullet_box(slide, left, top, width, height, title, bullets):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = GRAY_BORDER

    title_box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.10),
                                         width - Inches(0.36), Inches(0.35))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_run_style(run, 15, bold=True, color=TITLE_BLUE)

    content = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.46),
                                       width - Inches(0.36), height - Inches(0.56))
    tf = content.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.bullet = True
        p.space_after = Pt(6)
        for run in p.runs:
            set_run_style(run, 12, color=TEXT_DARK)
    return shape


def add_header(slide, title, subtitle):
    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.05)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = TITLE_BLUE
    banner.line.color.rgb = TITLE_BLUE

    add_textbox(slide, Inches(0.45), Inches(0.16), Inches(8.8), Inches(0.34),
                title, size=24, bold=True, color=RGBColor(255, 255, 255))
    add_textbox(slide, Inches(0.45), Inches(0.54), Inches(8.8), Inches(0.20),
                subtitle, size=11, color=RGBColor(230, 238, 248))

    add_textbox(
        slide, Inches(9.3), Inches(0.18), Inches(3.55), Inches(0.50),
        "Edite os campos: nome do supervisor, periodo e datas.",
        size=10, color=RGBColor(230, 238, 248), align=PP_ALIGN.RIGHT
    )


def add_card(slide, left, top, width, height, title):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = GRAY_BORDER
    add_textbox(slide, left + Inches(0.16), top + Inches(0.10), width - Inches(0.32),
                Inches(0.28), title, size=15, bold=True, color=TITLE_BLUE)
    return shape


def add_before_after_table(slide, left, top, width, height):
    add_card(slide, left, top, width, height, "Antes x Agora")
    col_gap = Inches(0.10)
    col_w = (width - Inches(0.48) - col_gap) / 2
    col1_left = left + Inches(0.16)
    col2_left = col1_left + col_w + col_gap
    head_top = top + Inches(0.46)
    row_h = Inches(0.60)

    add_textbox(slide, col1_left, head_top, col_w, Inches(0.32), "Antes",
                size=12, bold=True, color=TEXT_DARK, fill=RGBColor(255, 237, 237), margin=0.06)
    add_textbox(slide, col2_left, head_top, col_w, Inches(0.32), "Agora",
                size=12, bold=True, color=TEXT_DARK, fill=RGBColor(232, 245, 233), margin=0.06)

    left_items = [
        "Reagia de forma mais defensiva diante de negativas.",
        "Tinha dificuldade em ouvir feedback sem se justificar.",
        "Mostrava incomodo visivel em situacoes de contrariedade.",
        "Foco inicial maior na reacao do que na solucao.",
    ]
    right_items = [
        "Demonstra mais pausa antes de responder.",
        "Escuta com mais abertura e reduz justificativas.",
        "Mantem postura mais equilibrada e profissional.",
        "Apresenta maior foco em alternativas e solucao.",
    ]

    for idx, (item_left, item_right) in enumerate(zip(left_items, right_items)):
        row_top = head_top + Inches(0.38) + row_h * idx
        add_textbox(slide, col1_left, row_top, col_w, Inches(0.52), item_left,
                    size=11, color=TEXT_DARK, fill=RGBColor(252, 248, 248), margin=0.07)
        add_textbox(slide, col2_left, row_top, col_w, Inches(0.52), item_right,
                    size=11, color=TEXT_DARK, fill=RGBColor(246, 252, 247), margin=0.07)


def add_status_pill(slide, left, top, width, text, color):
    pill = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.28)
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    pill.line.color.rgb = color
    add_textbox(slide, left, top + Inches(0.01), width, Inches(0.20), text,
                size=10, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER, margin=0.03)


def add_semaforo_table(slide, left, top, width, height):
    add_card(slide, left, top, width, height, "Semaforo de comportamentos")
    rows = [
        ("Receber feedback corretivo", "Antes", RED, "Atual", YELLOW),
        ("Ouvir uma negativa", "Antes", RED, "Atual", GREEN),
        ("Ser contrariado em reuniao", "Antes", RED, "Atual", YELLOW),
        ("Responder sem justificativa excessiva", "Antes", RED, "Atual", YELLOW),
    ]
    col1 = left + Inches(0.16)
    col2 = left + Inches(2.95)
    col3 = left + Inches(3.80)
    body_top = top + Inches(0.52)

    add_textbox(slide, col1, body_top - Inches(0.20), Inches(2.65), Inches(0.24),
                "Situacao observada", size=10, bold=True, color=TEXT_MUTED)
    add_textbox(slide, col2, body_top - Inches(0.20), Inches(0.70), Inches(0.24),
                "Antes", size=10, bold=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, col3, body_top - Inches(0.20), Inches(0.72), Inches(0.24),
                "Atual", size=10, bold=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    for idx, (label, _, before_color, _, current_color) in enumerate(rows):
        row_top = body_top + Inches(0.56) * idx
        add_textbox(slide, col1, row_top, Inches(2.55), Inches(0.38), label, size=11, color=TEXT_DARK, margin=0.02)
        add_status_pill(slide, col2, row_top + Inches(0.04), Inches(0.58), "Risco", before_color)
        add_status_pill(slide, col3, row_top + Inches(0.04), Inches(0.62),
                        "Evolucao" if current_color == YELLOW else "Adequado", current_color)

    add_textbox(
        slide, left + Inches(0.16), top + height - Inches(0.42), width - Inches(0.32), Inches(0.24),
        "Legenda: vermelho = reacao inadequada | amarelo = parcial/em evolucao | verde = adequado",
        size=9, color=TEXT_MUTED
    )


def add_message_bar(slide, left, top, width):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.82)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE_LIGHT
    shape.line.color.rgb = RGBColor(185, 214, 255)
    add_textbox(slide, left + Inches(0.16), top + Inches(0.10), width - Inches(0.32), Inches(0.22),
                "Mensagem-chave", size=12, bold=True, color=TITLE_BLUE)
    add_textbox(
        slide, left + Inches(0.16), top + Inches(0.32), width - Inches(0.32), Inches(0.34),
        "Ha evolucao perceptivel no autocontrole, com avancos na escuta, menor impulsividade e comunicacao mais equilibrada diante de negativas e feedbacks.",
        size=11, color=TEXT_DARK
    )


def add_chart(slide, left, top, width, height):
    chart_data = CategoryChartData()
    chart_data.categories = ["Sem. 1", "Sem. 2", "Sem. 3", "Sem. 4"]
    chart_data.add_series("Controle emocional", (2, 3, 3, 4))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    ).chart
    chart.has_legend = False
    chart.value_axis.maximum_scale = 5
    chart.value_axis.minimum_scale = 0
    chart.value_axis.major_unit = 1
    chart.value_axis.has_major_gridlines = True
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = TITLE_BLUE
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = "Nota semanal de controle emocional (1 a 5)"
    for p in chart.chart_title.text_frame.paragraphs:
        for run in p.runs:
            set_run_style(run, 12, bold=True, color=TEXT_DARK)


def add_timeline(slide, left, top, width):
    add_textbox(slide, left, top, width, Inches(0.28), "Linha do tempo de evolucao",
                size=15, bold=True, color=TITLE_BLUE)
    line_top = top + Inches(0.58)
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, left + Inches(0.25), line_top, width - Inches(0.50), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(189, 199, 210)
    line.line.color.rgb = RGBColor(189, 199, 210)

    milestones = [
        ("Inicio", "Reacao mais defensiva e maior dificuldade com contrariedades."),
        ("Meio", "Passou a ouvir mais, apesar de ainda justificar em alguns momentos."),
        ("Atual", "Responde com maior equilibrio, mais escuta e melhor foco em solucao."),
    ]
    xs = [left + Inches(0.35), left + Inches(2.55), left + Inches(4.80)]
    for (label, desc), x in zip(milestones, xs):
        circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, line_top - Inches(0.12), Inches(0.28), Inches(0.28))
        circ.fill.solid()
        circ.fill.fore_color.rgb = TITLE_BLUE
        circ.line.color.rgb = TITLE_BLUE
        add_textbox(slide, x - Inches(0.05), line_top + Inches(0.20), Inches(1.00), Inches(0.22),
                    label, size=11, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
        add_textbox(slide, x - Inches(0.40), line_top + Inches(0.44), Inches(1.80), Inches(0.78),
                    desc, size=10, color=TEXT_DARK, align=PP_ALIGN.CENTER, margin=0.03)


def build_presentation(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    slide1 = prs.slides.add_slide(blank)
    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = BG_LIGHT
    add_header(
        slide1,
        "Evolucao do Controle Emocional do Supervisor",
        "Acompanhamento comportamental com base em situacoes reais de contrariedade, negativas e feedbacks.",
    )
    add_textbox(slide1, Inches(0.45), Inches(1.22), Inches(4.2), Inches(0.22),
                "Supervisor: [nome]   |   Periodo avaliado: [mes/ano a mes/ano]",
                size=11, color=TEXT_MUTED)
    add_before_after_table(slide1, Inches(0.45), Inches(1.60), Inches(6.35), Inches(4.35))
    add_semaforo_table(slide1, Inches(7.05), Inches(1.60), Inches(5.82), Inches(4.35))
    add_message_bar(slide1, Inches(0.45), Inches(6.10), Inches(12.42))

    slide2 = prs.slides.add_slide(blank)
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = BG_LIGHT
    add_header(
        slide2,
        "Evidencias de Evolucao",
        "Resumo visual para apresentacao de acompanhamento e desenvolvimento.",
    )
    add_chart(slide2, Inches(0.55), Inches(1.45), Inches(5.75), Inches(3.40))
    add_bullet_box(
        slide2, Inches(6.55), Inches(1.45), Inches(6.20), Inches(2.00), "Evidencias observadas",
        [
            "Em situacoes recentes de negativa, apresentou resposta mais contida e respeitosa.",
            "Demonstrou maior capacidade de ouvir feedback antes de responder.",
            "Reduziu a necessidade de justificativas longas em conversas de alinhamento.",
            "Passou a direcionar a conversa com mais foco em alternativas e solucao.",
        ]
    )
    add_bullet_box(
        slide2, Inches(6.55), Inches(3.65), Inches(6.20), Inches(1.90), "Leitura executiva",
        [
            "A evolucao e percebida pela reducao de reacoes impulsivas.",
            "Ha maior abertura para escuta e menor defensividade.",
            "O comportamento atual ja sustenta conversas mais produtivas e profissionais.",
        ]
    )
    add_timeline(slide2, Inches(0.55), Inches(5.15), Inches(6.15))
    add_bullet_box(
        slide2, Inches(7.05), Inches(5.15), Inches(5.70), Inches(1.55), "Fala de apoio",
        [
            "A comparacao entre o inicio do acompanhamento e o momento atual mostra melhora na escuta, no controle da reacao emocional e na reducao de justificativas excessivas.",
        ]
    )

    prs.save(output_path)


if __name__ == "__main__":
    build_presentation("/workspace/PDI_evolucao_supervisor.pptx")
